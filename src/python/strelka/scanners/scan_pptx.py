import io
import zipfile

from pptx import Presentation

from strelka import strelka


class ScanPptx(strelka.Scanner):
    """Collects metadata and extracts text from pptx files.

    Options:
        extract_text: Boolean that determines if document text should be
            extracted as a child file.
            Defaults to False.
    """

    def scan(self, data, file, options, expire_at):
        extract_text = options.get('extract_text', False)
        with io.BytesIO(data) as pptx_io:

            try:
                pptx_doc = Presentation(pptx_io)
                self.event['author'] = pptx_doc.core_properties.author
                self.event['category'] = pptx_doc.core_properties.category
                self.event['comments'] = pptx_doc.core_properties.comments
                self.event['content_status'] = pptx_doc.core_properties.content_status
                if pptx_doc.core_properties.created is not None:
                    self.event['created'] = int(pptx_doc.core_properties.created.strftime('%s'))
                self.event['identifier'] = pptx_doc.core_properties.identifier
                self.event['keywords'] = pptx_doc.core_properties.keywords
                self.event['language'] = pptx_doc.core_properties.language
                self.event['last_modified_by'] = pptx_doc.core_properties.last_modified_by
                if pptx_doc.core_properties.last_printed is not None:
                    self.event['last_printed'] = int(pptx_doc.core_properties.last_printed.strftime('%s'))
                if pptx_doc.core_properties.modified is not None:
                    self.event['modified'] = int(pptx_doc.core_properties.modified.strftime('%s'))
                self.event['revision'] = pptx_doc.core_properties.revision
                self.event['subject'] = pptx_doc.core_properties.subject
                self.event['title'] = pptx_doc.core_properties.title
                self.event['version'] = pptx_doc.core_properties.version
                self.event['slide_count'] = len(pptx_doc.slides)
                self.event['word_count'] = 0
                self.event['image_count'] = 0

                # Single pass: collect text, count words/images, extract hyperlinks
                extracted_text = [] if extract_text else None

                for slide in pptx_doc.slides:
                    for shape in slide.shapes:
                        # Count images
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            self.event['image_count'] += 1

                        # Process text frames
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                # Collect text for extraction
                                if extract_text and para.text:
                                    extracted_text.append(para.text)

                                # Count words
                                for run in para.runs:
                                    text = run.text.strip()
                                    if text:
                                        self.event['word_count'] += len(text.split())

                        # Extract hyperlinks
                        if hasattr(shape, 'click_action') and shape.click_action:
                            if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                                self.event.setdefault('hyperlinks', []).append(
                                    shape.click_action.hyperlink.address
                                )

                # Upload extracted text as single batch
                if extract_text and extracted_text:
                    extract_file = strelka.File(
                        name='text',
                        source=self.name,
                    )

                    text_content = '\n'.join(extracted_text)
                    for c in strelka.chunk_string(text_content):
                        self.upload_to_coordinator(
                            extract_file.pointer,
                            c,
                            expire_at,
                        )

                    self.files.append(extract_file)

            except ValueError:
                self.flags.append('value_error')
            except zipfile.BadZipFile:
                self.flags.append('bad_zip')
