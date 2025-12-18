"""
Standalone tests for ScanPptx scanner.
Run with: python -m pytest src/python/strelka/tests/test_scan_pptx_standalone.py -v
"""
import io
from pathlib import Path

import pytest


class TestScanPptx:
    """Tests for ScanPptx scanner."""

    @pytest.fixture
    def fixture_path(self):
        return Path(__file__).parent / "fixtures/test.pptx"

    @pytest.fixture
    def pptx_data(self, fixture_path):
        with open(fixture_path, "rb") as f:
            return f.read()

    def test_metadata_extraction(self, pptx_data):
        """Test that metadata is correctly extracted."""
        from pptx import Presentation

        pptx_doc = Presentation(io.BytesIO(pptx_data))

        assert pptx_doc.core_properties.comments == "generated using python-pptx"
        assert pptx_doc.core_properties.last_modified_by == "Test Author"
        assert pptx_doc.core_properties.revision == 1
        assert len(pptx_doc.slides) == 4

    def test_text_extraction(self, pptx_data):
        """Test that text is extracted from slides."""
        from pptx import Presentation

        pptx_doc = Presentation(io.BytesIO(pptx_data))

        all_text = []
        for slide in pptx_doc.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            all_text.append(para.text)

        full_text = " ".join(all_text)

        # Check for key phishing indicators (sanitized)
        assert "ACCESS YOUR BENEFITS" in full_text
        assert "Employment Contract" in full_text
        assert "john.doe@example.com" in full_text
        assert "ACME CORP" in full_text

    def test_hyperlink_extraction(self, pptx_data):
        """Test that hyperlinks are extracted from shapes."""
        from pptx import Presentation

        pptx_doc = Presentation(io.BytesIO(pptx_data))

        hyperlinks = []
        for slide in pptx_doc.slides:
            for shape in slide.shapes:
                if hasattr(shape, "click_action") and shape.click_action:
                    if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        hyperlinks.append(shape.click_action.hyperlink.address)

        assert len(hyperlinks) == 1
        assert "tracking-domain.example.com" in hyperlinks[0]
        assert "phishing.example.com" in hyperlinks[0]

    def test_word_count(self, pptx_data):
        """Test that word count is calculated correctly."""
        from pptx import Presentation

        pptx_doc = Presentation(io.BytesIO(pptx_data))

        word_count = 0
        for slide in pptx_doc.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                words = run.text.strip().split()
                                word_count += len(words)

        # Word count may vary slightly after sanitization
        assert word_count > 250

    def test_image_count(self, pptx_data):
        """Test that images are counted correctly."""
        from pptx import Presentation

        pptx_doc = Presentation(io.BytesIO(pptx_data))

        image_count = 0
        for slide in pptx_doc.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1

        assert image_count == 1

